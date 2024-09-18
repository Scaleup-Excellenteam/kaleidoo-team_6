import mimetypes
import pdfplumber
import docx
import pytesseract
from PIL import Image
import re
from pptx import Presentation
from io import BytesIO
from langdetect import detect
import chardet
import ftfy
from unidecode import unidecode

def read_file_with_encoding(file_path):
    """
    Tries to read the file using different encodings, starting with UTF-8 and falling back to ISO-8859-8.
    
    Args:
        file_path (str): Path to the file.
    
    Returns:
        str: The content of the file with correct encoding.
    """
    encodings = ['utf-8', 'iso-8859-8', 'windows-1255']  # Try common encodings for Hebrew text
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read(), encoding
        except UnicodeDecodeError:
            continue
    raise UnicodeDecodeError(f"Unable to decode file with the provided encodings: {encodings}")



def fix_unicode_text_if_needed(text, detected_language):
    """
    Fixes text encoding issues and applies necessary transformations.
    
    Args:
        text (str): The input text to be fixed.
        detected_language (str): The detected language of the text.
    
    Returns:
        str: The fixed text based on the language.
    """
    if detected_language != 'he':
        return unidecode(text)
    return ftfy.fix_text(text)

def clean_text(text):
    """
    clean_text(text: str) -> str
    Cleans and formats the input text by performing the following operations:
    1. Removes any space before punctuation marks such as `.`, `?`, `!`, `,`, `"` .
    2. Replaces multiple spaces with a single space and trims any leading or trailing whitespace.
    3. Replaces multiple consecutive newline characters with a single newline.
    
    Args:
        text (str): The input text to be cleaned and formatted.
    
    Returns:
        str: The cleaned text with spaces and newlines properly formatted.
    """
    text = re.sub(r'\s([?.!,"](?:\s|$))', r'\1', text)
    text = re.sub(r'\s+', ' ', text).strip()
    text = re.sub(r'[\r\n]+', '\n', text)  
    return text



def extract_text_from_handwriting(image):
    """
    Utilizes pytesseract to extract text from handwriting in the given image.
    
    Args:
        image (PIL.Image): The image from which text needs to be extracted.
    
    Returns:
        str: The extracted text from handwriting.
    """
    return pytesseract.image_to_string(image, config='--oem 1 --psm 3')



def extract_tables_from_pdf(pdf):
    """
    Extracts tables from a PDF file using pdfplumber.
    
    Args:
        pdf (pdfplumber.PDF): The PDF object to extract tables from.
    
    Returns:
        str: The extracted tables as a string.
    """
    tables_text = ""
    for page in pdf.pages:
        tables = page.extract_tables()
        for table in tables:
            for row in table:
                row = [cell if cell is not None else '' for cell in row]
                tables_text += "\t".join(row) + "\n"
    return tables_text


def extract_tables_from_docx(doc):
    """
    Extracts tables from a Word document using python-docx.
    
    Args:
        doc (docx.Document): The Word document object to extract tables from.
    
    Returns:
        str: The extracted tables as a string.
    """
    tables_text = ""
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text for cell in row.cells]
            tables_text += "\t".join(row_text) + "\n"
    return tables_text

def extract_text_from_file(file_path):
    """
    Extracts text and tables from a file based on its MIME type. 
    Handles plain text, PDF, Word documents, PowerPoint presentations, and images.
    
    Args:
        file_path (str): Path to the file to extract text and tables from.
    
    Returns:
        str: The extracted text and tables from the file.
    """
    mime_type, _ = mimetypes.guess_type(file_path)

    if mime_type == 'text/plain':
        text, encoding_used = read_file_with_encoding(file_path)
 

    elif mime_type == 'application/pdf':
        with pdfplumber.open(file_path) as pdf:
            text = ""
            for page in pdf.pages:
                text += page.extract_text()
            text += "\n" + extract_tables_from_pdf(pdf)

    elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        text += "\n" + extract_tables_from_docx(doc)
        # Handling images in Word
        for rel in doc.part.rels:
            if "image" in doc.part.rels[rel].target_ref:
                image_data = doc.part.rels[rel].target_part.blob
                image = Image.open(BytesIO(image_data))
                text += "\n" + extract_text_from_handwriting(image)

    elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image = Image.open(BytesIO(shape.image.blob))
                    text += "\n" + extract_text_from_handwriting(image)

    elif mime_type.startswith('image/'):
        image = Image.open(file_path)
        text = extract_text_from_handwriting(image)

    else:
        raise ValueError("Unsupported file type")

    try:
        detected_language = detect(text)
    except:
        detected_language = 'unknown'
        print("Unable to detect language")

    return text, detected_language

def save_text_to_file(file_path, fixed_dest_folder):
    """
    Saves the extracted text to a file.
    
    Args:
        text (str): The text to be saved.
        file_path (str): The path to save the text to.
    """
    context, language = extract_text_from_file(file_path)
    file_name = file_path.split("\\")[-1].split(".")[0]
    fixed_text = fix_unicode_text_if_needed(context, language)
    write_path = f"{fixed_dest_folder}/{file_name}_fixed.txt"
    with open(write_path, "w", encoding="utf-8") as file:
        file.write(fixed_text)


